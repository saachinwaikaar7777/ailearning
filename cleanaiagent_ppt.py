import os
import json
import re
import urllib.parse
import requests
import warnings
import time
from datetime import datetime

warnings.filterwarnings("ignore", category=DeprecationWarning)

from dotenv import load_dotenv

# --- LLM Imports (Uncomment the ChatOllama import if you switch to Ollama) ---
from langchain_google_genai import ChatGoogleGenerativeAI
#from langchain_community.chat_models import ChatOllama

from langchain_core.messages import SystemMessage, HumanMessage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

# ──────────────────────────────────────────────────────────────
# Step 0: Choose Your Engine (Gemini vs. Ollama)
# ──────────────────────────────────────────────────────────────

# --- OPTION A: Gemini API (Default Cloud) ---
llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash-lite", temperature=0, max_retries=5)

# --- OPTION B: Ollama (Local/Offline) ---
# To use Ollama, run: pip install langchain-community
# Make sure Ollama is running on your machine and you have pulled a model (e.g., ollama pull llama3)
#lm = ChatOllama(model="llama3:latest", temperature=0)


# ──────────────────────────────────────────────────────────────
# Global Configuration & Design System Presets
# ──────────────────────────────────────────────────────────────

DECK_ARCHETYPES = {
    "1": {
        "name": "Sales Deck",
        "primary": RGBColor(0x4F, 0x46, 0xE5),   # Modern Indigo 600
        "body": RGBColor(0x1F, 0x29, 0x37),      # Dark Charcoal Gray 800
        "accent": RGBColor(0xF5, 0x9E, 0x0B),    # Attention Amber 500
        "bg": RGBColor(0xFF, 0xFF, 0xFF),        # Pure White
        "card": RGBColor(0xF9, 0xFA, 0xFB),      # Cool Soft Gray 50
        "font_title": "Trebuchet MS",            # Dynamic/punchy title
        "font_body": "Arial",
        "tone": "highly persuasive, high-energy, value-proposition focused, highlighting client pain points, solutions, and key metrics."
    },
    "2": {
        "name": "Financial Report",
        "primary": RGBColor(0x0F, 0x17, 0x2A),   # Slate 900
        "body": RGBColor(0x33, 0x41, 0x55),      # Slate 700
        "accent": RGBColor(0x0D, 0x94, 0x88),    # Teal 600 (Represents growth)
        "bg": RGBColor(0xFC, 0xFD, 0xFD),        # Ultra-clean off-white
        "card": RGBColor(0xF1, 0xF5, 0xF9),      # Cool Slate Blue Card 100
        "font_title": "Georgia",                 # Corporate Editorial Serif
        "font_body": "Calibri",                  # Highly readable
        "tone": "conservative, precise, metric-driven, structured around quarters, percentages, and financial performance markers."
    },
    "3": {
        "name": "Academic Slides",
        "primary": RGBColor(0x1E, 0x3A, 0x8A),   # Deep Intellectual Blue 900
        "body": RGBColor(0x1F, 0x29, 0x37),      # Gray 800
        "accent": RGBColor(0x99, 0x1B, 0x1B),    # Ivy Crimson Red 800
        "bg": RGBColor(0xFF, 0xFF, 0xF8),        # Warm Scholar Ivory
        "card": RGBColor(0xF5, 0xF5, 0xF0),      # Soft muted beige card
        "font_title": "Times New Roman",         # Formal Editorial Serif
        "font_body": "Calibri",
        "tone": "highly explanatory, analytical, structured around research methodology, hypotheses, and conclusions."
    },
    "4": {
        "name": "Default Minimalist",
        "primary": RGBColor(0x0F, 0x17, 0x2A),   # Slate 900
        "body": RGBColor(0x33, 0x41, 0x55),      # Slate 700
        "accent": RGBColor(0x25, 0x63, 0xEB),    # Royal Blue 600
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "card": RGBColor(0xF8, 0xFA, 0xFC),      # Slate 50
        "font_title": "Georgia",
        "font_body": "Arial",
        "tone": "neutral, objective, clean, prioritizing minimalist structural layouts."
    }
}

USER_CONFIG = {
    "template_path": "",
    "branding": "",
    "deck_type": "4" # Default Archetype Key
}

# ──────────────────────────────────────────────────────────────
# Token-Optimized Internal Helpers
# ──────────────────────────────────────────────────────────────

def fetch_wikipedia_research(query: str) -> str:
    """
    Executes a direct programmatic Wikipedia API query in Python.
    Consumes 0 LLM tokens, saving massive search-loop overhead.
    """
    url = f"https://en.wikipedia.org/w/api.php?action=query&list=search&srsearch={urllib.parse.quote(query)}&utf8=&format=json"
    try:
        response = requests.get(url, timeout=10).json()
        snippets = [item['snippet'] for item in response['query']['search'][:5]]
        clean_snippets = [re.sub('<[^<]+>', '', s) for s in snippets]
        if not clean_snippets:
            return ""
        return "\n".join(f"- {s}" for s in clean_snippets)
    except Exception:
        return ""


def invoke_with_retry(llm_instance, inputs, max_attempts=5):
    """
    Invokes the LLM with dynamic 429 rate-limit backoff recovery.
    No-op if running local Ollama engines.
    """
    # Auto-detect if using local Ollama (which does not have cloud rate-limits)
    if "Ollama" in type(llm_instance).__name__:
        return llm_instance.invoke(inputs)

    for attempt in range(max_attempts):
        try:
            return llm_instance.invoke(inputs)
        except Exception as e:
            err_msg = str(e)
            if "429" in err_msg or "RESOURCE_EXHAUSTED" in err_msg:
                match = re.search(r"retry in ([\d\.]+)s", err_msg)
                sleep_time = float(match.group(1)) if match else (2 ** attempt) * 5
                print(f"\n⚠️  [Rate Limit Active] Gemini Free Tier quota reached.")
                print(f"    Waiting for {sleep_time:.2f} seconds before continuing (Attempt {attempt + 1}/{max_attempts})...")
                time.sleep(sleep_time + 1.5)
            else:
                raise e
    return llm_instance.invoke(inputs)


def compile_presentation(slides_json: str) -> str:
    """
    Direct Python PowerPoint Compilation Engine.
    Bypasses LLM planning graphs entirely to render pixels programmatically.
    """
    try:
        triple_ticks = "`" * 3
        formatting_pattern = rf"{triple_ticks}(?:json)?\n?(.*?)\n?{triple_ticks}"
        clean_json = re.sub(formatting_pattern, r"\1", slides_json, flags=re.DOTALL).strip()
        slides_data = json.loads(clean_json)
    except json.JSONDecodeError as e:
        return f"ERROR: Invalid JSON provided – {e}. Please ensure output is valid JSON."

    output_folder = "output_PPT"
    os.makedirs(output_folder, exist_ok=True)
    
    use_template = bool(USER_CONFIG["template_path"] and os.path.exists(USER_CONFIG["template_path"]))
    archetype_key = USER_CONFIG.get("deck_type", "4")
    style = DECK_ARCHETYPES.get(archetype_key, DECK_ARCHETYPES["4"])
    
    COLOR_PRIMARY = style["primary"]
    COLOR_BODY    = style["body"]
    COLOR_ACCENT  = style["accent"]
    COLOR_CARD_BG = style["card"]
    COLOR_BG      = style["bg"]
    FONT_TITLE    = style["font_title"]
    FONT_BODY     = style["font_body"]

    if use_template:
        prs = Presentation(USER_CONFIG["template_path"])
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]
    else:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
    def add_card(slide, left, top, width, height, fill_color=COLOR_CARD_BG):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        shape.line.width = Pt(1)
        return shape

    def add_textbox(slide, text, left, top, width, height, font_size=16, font_color=COLOR_BODY, bold=False, font_name=FONT_BODY, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.bold = bold
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        return txBox

    def add_bullet_textbox(slide, text, left, top, width, height, font_size=14, font_name=FONT_BODY):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        
        run_bullet = p.add_run()
        run_bullet.text = "•  "
        run_bullet.font.name = font_name
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = COLOR_ACCENT
        run_bullet.font.bold = True

        clean_text = text.lstrip("•").strip()

        if ":" in clean_text:
            anchor, description = clean_text.split(":", 1)
            run_anchor = p.add_run()
            run_anchor.text = anchor.strip() + ": "
            run_anchor.font.name = font_name
            run_anchor.font.bold = True
            run_anchor.font.size = Pt(font_size)
            run_anchor.font.color.rgb = COLOR_PRIMARY
            
            run_desc = p.add_run()
            run_desc.text = description.strip()
            run_desc.font.name = font_name
            run_desc.font.bold = False
            run_desc.font.size = Pt(font_size)
            run_desc.font.color.rgb = COLOR_BODY
        else:
            run_text = p.add_run()
            run_text.text = clean_text
            run_text.font.name = font_name
            run_text.font.bold = False
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = COLOR_BODY
            
        return txBox

    def fetch_image(prompt, index):
        try:
            safe_prompt = urllib.parse.quote(prompt + " " + USER_CONFIG["branding"] + ", clean studio photography, premium look, no text")
            url = f"https://image.pollinations.ai/prompt/{safe_prompt}?width=1024&height=768&nologo=true"
            response = requests.get(url, timeout=12)
            if response.status_code == 200:
                img_path = os.path.join(output_folder, f"temp_img_{index}.jpg")
                with open(img_path, 'wb') as f:
                    f.write(response.content)
                return img_path
        except Exception as e:
            print(f"  [!] Image omitted for slide {index}: {e}")
        return None

    # Slide 0 Title Setup
    title_info = slides_data[0] if slides_data else {"title": "Strategic Overview", "bullets": ["Generated via AI Engine"]}
    
    if use_template:
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title: slide.shapes.title.text = title_info["title"]
        if len(slide.placeholders) > 1 and title_info.get("bullets"):
            slide.placeholders[1].text = title_info["bullets"][0]
    else:
        slide = prs.slides.add_slide(blank_layout)
        add_card(slide, 0, 0, 13.33, 7.5, fill_color=COLOR_BG)
        
        accent_bar = slide.shapes.add_shape(1, Inches(1.2), Inches(2.2), Inches(0.15), Inches(2.8))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = COLOR_ACCENT
        accent_bar.line.fill.background()

        add_textbox(slide, title_info["title"], left=1.6, top=2.1, width=10.5, height=1.8, font_size=42, font_color=COLOR_PRIMARY, bold=True, font_name=FONT_TITLE)
        if title_info.get("bullets"):
            add_textbox(slide, title_info["bullets"][0].upper(), left=1.6, top=4.1, width=10.5, height=0.6, font_size=13, font_color=COLOR_BODY, bold=True, font_name=FONT_BODY)

    # Content Slides rendering
    for idx, slide_info in enumerate(slides_data[1:], start=1):
        img_path = slide_info.get("image_prompt") and fetch_image(slide_info["image_prompt"], idx)

        if use_template:
            slide = prs.slides.add_slide(content_layout)
            if slide.shapes.title: slide.shapes.title.text = slide_info.get("title", "")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text_frame.text = "\n".join(slide_info.get("bullets", []))
            if img_path:
                slide.shapes.add_picture(img_path, Inches(8.5), Inches(2.0), width=Inches(4.0))
        else:
            slide = prs.slides.add_slide(blank_layout)
            add_card(slide, 0, 0, 13.33, 7.5, fill_color=COLOR_BG)

            add_textbox(slide, slide_info.get("title", ""), left=0.8, top=0.6, width=11.5, height=0.8, font_size=28, font_color=COLOR_PRIMARY, bold=True, font_name=FONT_TITLE)

            left_margin = 0.8
            available_width = 5.6 if img_path else 11.7
            
            add_card(slide, left_margin, 1.8, available_width, 4.8)

            current_y = 2.1
            bullets = slide_info.get("bullets", [])
            
            for bullet in bullets:
                add_bullet_textbox(slide, bullet, left=left_margin + 0.4, top=current_y, width=available_width - 0.8, height=0.8, font_size=14, font_name=FONT_BODY)
                current_y += 1.0
                if current_y > 6.0: 
                    break

            if img_path:
                slide.shapes.add_picture(img_path, Inches(6.8), Inches(1.8), width=Inches(5.7), height=Inches(4.8))

    # Cleanup temporary images
    for idx in range(1, len(slides_data)):
        tmp_img = os.path.join(output_folder, f"temp_img_{idx}.jpg")
        if os.path.exists(tmp_img): os.remove(tmp_img)

    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    output_path = os.path.join(output_folder, f"{timestamp}.pptx")
    prs.save(output_path)
    return f"SUCCESS: Presentation saved to '{output_path}' with {len(prs.slides)} slides."


# ──────────────────────────────────────────────────────────────
# Step 3: Run Orchestration Framework (Optimized Single LLM Call)
# ──────────────────────────────────────────────────────────────

def run_agents(topic: str):
    print(f"\n{'='*60}")
    print(f"  Generating {DECK_ARCHETYPES[USER_CONFIG['deck_type']]['name']} for: {topic}")
    print(f"{'='*60}")

    # Stage 1: Fast Programmatic Wikipedia Hook (0 LLM Tokens)
    print("\n[1/2] 🔍 Gathering Strategic Domain Research (0 tokens used)...")
    research_output = fetch_wikipedia_research(topic)

    # Stage 2: Direct Single-Turn Synthesis (Saves up to 70% of total LLM cost!)
    print("\n[2/2] ✍️ Generating Slide Architecture & Rendering Graphics...")
    archetype = DECK_ARCHETYPES.get(USER_CONFIG["deck_type"], DECK_ARCHETYPES["4"])
    type_name = archetype["name"]
    tone_focus = archetype["tone"]
    branding = USER_CONFIG["branding"]
    
    system_prompt = (
        f"You are an Expert Designer and Copywriter specializing in structural layouts for a {type_name}.\n"
        "Convert the provided raw research insights into an elegant, machine-readable JSON array representing PowerPoint slides.\n\n"
        f"DESIGN THEME STYLE PROFILE:\n- Tone: {tone_focus}\n"
        f"- Branding Preferences: '{branding if branding else 'Minimalist layout, clean contrast'}'\n\n"
        "CRITICAL STRUCTURAL RULES:\n"
        "1. Output ONLY a raw, parser-safe JSON array. Do not append Markdown blocks, conversational prefaces, or wrappers.\n"
        "2. Slide 0 must be your Master Title Slide: {{\"title\": \"Main Presentation Title\", \"bullets\": [\"Sub-insight or subtitle contextualizing value\"], \"image_prompt\": \"\"}}\n"
        "3. Subsequent slides (4-6 total elements) must map to this exact structural schema:\n"
        "   {{\"title\": \"Actionable Slide Header\", \"bullets\": [\"Anchor: Concise descriptive context, not exceed 12 words.\"], \"image_prompt\": \"Descriptive, professional visual concept, high-contrast studio lighting, 8k resolution, photorealistic. NO text overlays.\"}}\n"
        "4. IMPORTANT: Format every bullet point using an Anchor Word followed by a colon (e.g., 'Market Growth: The industry scaling exponentially.'). This enforces professional visual layout styling.\n"
        "5. Maintain absolute structural uniformity across every slide object inside the array."
    )
    
    user_prompt = (
        f"Topic: {topic}\n\n"
        f"Wikipedia Research Facts:\n"
        f"{research_output if research_output else 'None found on Wikipedia. Rely entirely on your internal expertise.'}"
    )
    
    messages = [
        SystemMessage(content=system_prompt),
        HumanMessage(content=user_prompt)
    ]
    
    # Executes the only LLM call of the entire run
    writer_result = invoke_with_retry(llm, messages)
    slides_json = writer_result.content
    
    # Run the dynamic compilation pipeline directly in Python (0 LLM overhead)
    ppt_output = compile_presentation(slides_json)
    
    print("\n✅ Execution Framework Completed Successfully:")
    print(ppt_output)

if __name__ == "__main__":
    print(r"""
  __  __        _ _   _              _                    _       ____  ____ _____ 
 |  \/  |     | | | (_)            / \   __ _  ___ _ __ | |_    |  _ \|  _ \_   _|
 | |\/| | |  | | | |_| |  _____    / _ \ / _` |/ _ \ '_ \ __|   | |_) | |_) || |  
 | |  | | |_| | |  _| | |_____|  / ___ \ (_| |  __/ | | | |_    |  __/|  __/ | |  
 |_|  |_|\__,_|_|\__|_|         /_/   \_\__, |\___|_| |_|\__|   |_|   |_|    |_|  
                                         |___/                                    
    """)
    
    user_topic = input("\n👉 Enter the topic for your presentation: ").strip()
    if not user_topic:
        user_topic = "The Future of Artificial Intelligence"
        print(f"No topic provided. Defaulting to: {user_topic}")
    
    print("\n👉 Choose Presentation Type:")
    print("   [1] Sales Deck (Persuasive, Bold Accent, Modern)")
    print("   [2] Financial Report (Metric-driven, Conservative, Teal/Navy Slate)")
    print("   [3] Academic Slides (Exploratory, Rich Literature, Ivory Background)")
    print("   [4] Default Minimalist (Clean, Corporate Slate, Royal Blue)")
    deck_choice = input("   Selection (1-4): ").strip()
    if deck_choice in DECK_ARCHETYPES:
        USER_CONFIG["deck_type"] = deck_choice
        print(f"   [+] Configuration loaded: {DECK_ARCHETYPES[deck_choice]['name']}")
    else:
        USER_CONFIG["deck_type"] = "4"
        print("   [!] Invalid selection. Defaulting to Minimalist.")

    print("\n👉 Do you have a custom PowerPoint template? ")
    user_template = input("   (Enter the path to the .pptx file, or leave blank to use AI-generated design): ").strip()
    if user_template and os.path.exists(user_template):
        USER_CONFIG["template_path"] = user_template
        print(f"   [+] Valid template found: {user_template}")
    else:
        print("   [+] Using default AI structural canvas layout engine.")

    print("\n👉 Enter any specific branding/style guidelines.")
    user_branding = input("   (e.g., 'Minimalist design, corporate blue, futuristic vibe', or leave blank): ").strip()
    USER_CONFIG["branding"] = user_branding

    run_agents(user_topic)