"""
Multi-Agent PPT Creator
Pipeline: Research Agent → Writer Agent → PPT Agent
"""
import os
import json
import re
import warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

from dotenv import load_dotenv
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.tools import tool
from langgraph.prebuilt import create_react_agent

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash-lite", temperature=0)


# ──────────────────────────────────────────────────────────────
# Step 2: Tools
# ──────────────────────────────────────────────────────────────

@tool
def research_tool(query: str) -> str:
    """Useful for researching any topic and returning key facts."""
    return (
        f"Research findings on '{query}':\n"
        "1. Regular exercise reduces the risk of chronic diseases by up to 50%.\n"
        "2. Physical activity boosts productivity and focus at work.\n"
        "3. Active individuals earn ~10% more on average due to higher energy & discipline.\n"
        "4. Exercise reduces healthcare costs significantly over a lifetime.\n"
        "5. Mental health improvements from exercise reduce absenteeism."
    )


@tool
def create_ppt(slides_json: str) -> str:
    """
    Creates a PowerPoint file from a JSON list of slides.

    Expected JSON format:
    [
      {"title": "Slide Title", "bullets": ["Point 1", "Point 2", ...]},
      ...
    ]

    Returns the path to the saved .pptx file.
    """
    # ── Parse input ───────────────────────────────────────────
    try:
        slides_data = json.loads(slides_json)
    except json.JSONDecodeError as e:
        return f"ERROR: Invalid JSON provided – {e}"

    # ── Output folder ─────────────────────────────────────────
    output_folder = "output_PPT"                          # ← folder name
    os.makedirs(output_folder, exist_ok=True)             # ← creates if missing

    # ── Theme colours ─────────────────────────────────────────
    DARK_BG   = RGBColor(0x1E, 0x27, 0x61)   # navy
    ACCENT    = RGBColor(0xCA, 0xDC, 0xFC)   # ice blue
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xFF)   # near-white blue tint
    BODY_TEXT = RGBColor(0x1E, 0x27, 0x61)   # navy on light slides

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]   # fully blank

    def add_rect(slide, l, t, w, h, fill_color):
        shape = slide.shapes.add_shape(
            1, Inches(l), Inches(t), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def add_textbox(slide, text, l, t, w, h, bold=False, size=18,
                    color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(
            Inches(l), Inches(t), Inches(w), Inches(h)
        )
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold  = bold
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        return txBox

    # ── Slide 0: Title slide (dark) ───────────────────────────
    title_info = slides_data[0] if slides_data else {"title": "Presentation", "bullets": []}
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    add_rect(slide, 0, 5.8, 13.33, 1.7, RGBColor(0x16, 0x1D, 0x4E))

    add_textbox(slide, title_info["title"],
                l=1.0, t=2.2, w=11.0, h=1.8,
                bold=True, size=40, color=WHITE, align=PP_ALIGN.CENTER)

    subtitle = title_info["bullets"][0] if title_info.get("bullets") else ""
    if subtitle:
        add_textbox(slide, subtitle,
                    l=1.5, t=4.2, w=10.0, h=0.9,
                    bold=False, size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    # ── Content slides (light) ────────────────────────────────
    for slide_info in slides_data[1:]:
        slide = prs.slides.add_slide(blank_layout)
        add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_BG)
        add_rect(slide, 0, 0, 13.33, 1.35, DARK_BG)

        add_textbox(slide, slide_info.get("title", ""),
                    l=0.4, t=0.15, w=12.5, h=1.0,
                    bold=True, size=28, color=WHITE)

        bullets = slide_info.get("bullets", [])
        y = 1.65
        for bullet in bullets:
            dot = slide.shapes.add_shape(
                1, Inches(0.45), Inches(y + 0.07), Inches(0.18), Inches(0.18)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = DARK_BG
            dot.line.fill.background()

            add_textbox(slide, bullet,
                        l=0.75, t=y, w=12.0, h=0.55,
                        bold=False, size=16, color=BODY_TEXT)
            y += 0.72
            if y > 6.8:
                break

    # ── Last slide: closing ───────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    add_textbox(slide, "Thank You",
                l=1.0, t=2.8, w=11.0, h=1.5,
                bold=True, size=44, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Questions & Discussion",
                l=2.0, t=4.5, w=9.0, h=0.8,
                bold=False, size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    # ── Save to output_PPT folder ─────────────────────────────
    output_path = os.path.join(output_folder, "output_presentation.pptx")  # ← updated path
    prs.save(output_path)
    return f"SUCCESS: Presentation saved to '{output_path}' with {len(prs.slides)} slides."


# ──────────────────────────────────────────────────────────────
# Step 3: Agents
# ──────────────────────────────────────────────────────────────

research_agent = create_react_agent(
    model=llm,
    tools=[research_tool],
    prompt=(
        "You are a research agent. Use the research_tool to gather "
        "detailed, factual information about the given topic. "
        "Return a thorough summary of your findings."
    )
)

writer_agent = create_react_agent(
    model=llm,
    tools=[],
    prompt=(
        "You are a presentation writer. Given research findings, produce "
        "a structured JSON array suitable for a PowerPoint presentation.\n\n"
        "RULES:\n"
        "- Output ONLY valid JSON — no markdown fences, no extra text.\n"
        "- First element must be the title slide: "
        '{"title": "<Presentation Title>", "bullets": ["<Subtitle or tagline>"]}\n'
        "- Follow with 4–6 content slides: "
        '{"title": "<Slide Title>", "bullets": ["<point>", "<point>", ...]}\n'
        "- Each slide should have 3–5 concise bullet points (max 12 words each).\n"
        "- Bullet text must be factual and directly from the research."
    )
)

ppt_agent = create_react_agent(
    model=llm,
    tools=[create_ppt],
    prompt=(
        "You are a PowerPoint creation agent. You receive a JSON array "
        "describing slide content. Call the create_ppt tool with that JSON "
        "to generate the presentation file. "
        "If the input is not valid JSON, ask for a corrected version."
    )
)


# ──────────────────────────────────────────────────────────────
# Step 4: Orchestration
# ──────────────────────────────────────────────────────────────

def run_agents(topic: str) -> str:
    print(f"\n{'='*55}")
    print(f"  Topic: {topic}")
    print(f"{'='*55}")

    # ── Stage 1: Research ─────────────────────────────────────
    print("\n[1/3] Research Agent running…")
    research_result = research_agent.invoke(
        {"messages": [{"role": "user", "content": topic}]}
    )
    research_output = research_result["messages"][-1].content
    print(f"Research Output:\n{research_output}")

    # ── Stage 2: Writer → structured JSON ────────────────────
    print("\n[2/3] Writer Agent structuring slides…")
    writer_prompt = (
        f"Convert the following research into a PowerPoint JSON structure:\n\n"
        f"{research_output}"
    )
    writer_result = writer_agent.invoke(
        {"messages": [{"role": "user", "content": writer_prompt}]}
    )
    slides_json = writer_result["messages"][-1].content

    # Strip accidental markdown fences if present
    slides_json = re.sub(r"```(?:json)?", "", slides_json).strip()
    print(f"Slides JSON:\n{slides_json}")

    # ── Stage 3: PPT Agent creates file ───────────────────────
    print("\n[3/3] PPT Agent creating presentation…")
    ppt_result = ppt_agent.invoke(
        {"messages": [{"role": "user", "content": slides_json}]}
    )
    ppt_output = ppt_result["messages"][-1].content
    print(f"PPT Agent Output: {ppt_output}")

    return ppt_output


# ──────────────────────────────────────────────────────────────
# Step 5: Entry point
# ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    result = run_agents("Impact of exercise on wealth")
    print(f"\n{'='*55}")
    print("DONE:", result)
    print(f"{'='*55}")