import os
import json
import re
import urllib.parse
import requests
import warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

from dotenv import load_dotenv
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.tools import tool
from langgraph.prebuilt import create_react_agent

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

# Ensure you have your GOOGLE_API_KEY in your .env file
llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash-lite", temperature=0)

# Global Configuration object to hold user inputs
USER_CONFIG = {
    "template_path": "",
    "branding": ""
}

# ──────────────────────────────────────────────────────────────
# Step 1: Tools
# ──────────────────────────────────────────────────────────────

@tool
def research_tool(query: str) -> str:
    """Useful for researching factual data about a given topic."""
    url = f"https://en.wikipedia.org/w/api.php?action=query&list=search&srsearch={urllib.parse.quote(query)}&utf8=&format=json"
    try:
        response = requests.get(url).json()
        snippets = [item['snippet'] for item in response['query']['search'][:5]]
        # Clean HTML tags from Wikipedia snippets
        clean_snippets = [re.sub('<[^<]+>', '', s) for s in snippets]
        
        if not clean_snippets:
            return f"No specific facts found for '{query}', use general knowledge."
            
        return f"Research facts for '{query}':\n" + "\n".join(f"- {s}" for s in clean_snippets)
    except Exception as e:
        return f"Could not connect to external search. Rely on internal knowledge for: {query}"


@tool
def create_ppt(slides_json: str) -> str:
    """
    Creates a PowerPoint file from a JSON list of slides.
    Downloads images dynamically based on 'image_prompt'.
    """
    # ── Parse and clean input ─────────────────────────────────
    try:
        # Strip potential markdown code blocks
        clean_json = re.sub(r"```(?:json)?\n?(.*?)\n?```", r"\1", slides_json, flags=re.DOTALL).strip()
        slides_data = json.loads(clean_json)
    except json.JSONDecodeError as e:
        return f"ERROR: Invalid JSON provided – {e}. Please return ONLY valid JSON."

    output_folder = "output_PPT"
    os.makedirs(output_folder, exist_ok=True)
    
    # Check if a valid custom template was provided
    use_template = bool(USER_CONFIG["template_path"] and os.path.exists(USER_CONFIG["template_path"]))
    
    if use_template:
        prs = Presentation(USER_CONFIG["template_path"])
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1] # standard Title + Content
    else:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        # Colors for default generation
        DARK_BG   = RGBColor(0x1E, 0x27, 0x61)
        WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xFF)
        BODY_TEXT = RGBColor(0x1E, 0x27, 0x61)

        def add_rect(slide, l, t, w, h, fill_color):
            shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            shape.line.fill.background()
            return shape

        def add_textbox(slide, text, l, t, w, h, bold=False, size=18, color=WHITE, align=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.bold  = bold
            run.font.size  = Pt(size)
            run.font.color.rgb = color
            return txBox

    # Helper function to fetch images using free Pollinations API
    def fetch_image(prompt, index):
        try:
            # High-quality AI image generation without API keys
            safe_prompt = urllib.parse.quote(prompt + " " + USER_CONFIG["branding"] + ", no text, highly detailed")
            url = f"https://image.pollinations.ai/prompt/{safe_prompt}?width=800&height=600&nologo=true"
            response = requests.get(url, timeout=15)
            if response.status_code == 200:
                img_path = os.path.join(output_folder, f"temp_img_{index}.jpg")
                with open(img_path, 'wb') as f:
                    f.write(response.content)
                return img_path
        except Exception as e:
            print(f"  [!] Failed to generate image for slide {index}: {e}")
        return None

    # ── Slide 0: Title slide ──────────────────────────────────
    title_info = slides_data[0] if slides_data else {"title": "Presentation", "bullets": []}
    
    if use_template:
        slide = prs.slides.add_slide(title_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title_info["title"]
        if len(slide.placeholders) > 1 and title_info.get("bullets"):
            slide.placeholders[1].text = title_info["bullets"][0]
    else:
        slide = prs.slides.add_slide(blank_layout)
        add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
        add_textbox(slide, title_info["title"], l=1.0, t=2.5, w=11.0, h=1.8, bold=True, size=40, align=PP_ALIGN.CENTER)
        if title_info.get("bullets"):
            add_textbox(slide, title_info["bullets"][0], l=1.5, t=4.2, w=10.0, h=0.9, size=20, align=PP_ALIGN.CENTER)

    # ── Content slides ────────────────────────────────────────
    for idx, slide_info in enumerate(slides_data[1:], start=1):
        # 1. Download image if a prompt exists
        img_path = None
        if slide_info.get("image_prompt"):
            img_path = fetch_image(slide_info["image_prompt"], idx)

        # 2. Add slide to PPT
        if use_template:
            slide = prs.slides.add_slide(content_layout)
            if slide.shapes.title:
                slide.shapes.title.text = slide_info.get("title", "")
            
            # Add bullets
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.text = "\n".join(slide_info.get("bullets", []))
            
            # Add image to top right corner if it exists
            if img_path:
                slide.shapes.add_picture(img_path, Inches(8.5), Inches(2.0), width=Inches(4.0))

        else:
            slide = prs.slides.add_slide(blank_layout)
            add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_BG)
            add_rect(slide, 0, 0, 13.33, 1.35, DARK_BG)

            add_textbox(slide, slide_info.get("title", ""), l=0.4, t=0.15, w=12.5, h=1.0, bold=True, size=28)

            y = 1.65
            # If we have an image, restrict text width to left half
            text_width = 6.0 if img_path else 12.0
            
            for bullet in slide_info.get("bullets", []):
                add_textbox(slide, "• " + bullet, l=0.5, t=y, w=text_width, h=0.55, size=16, color=BODY_TEXT)
                y += 0.8
                if y > 6.8:
                    break
            
            if img_path:
                # Add image on the right side
                slide.shapes.add_picture(img_path, Inches(6.8), Inches(1.8), width=Inches(5.8))

    # Clean up temporary images
    for idx in range(1, len(slides_data)):
        tmp_img = os.path.join(output_folder, f"temp_img_{idx}.jpg")
        if os.path.exists(tmp_img):
            os.remove(tmp_img)

    # ── Save Output ───────────────────────────────────────────
    output_path = os.path.join(output_folder, f"{title_info['title'][:10].replace(' ','_')}_presentation.pptx")
    prs.save(output_path)
    return f"SUCCESS: Presentation saved to '{output_path}' with {len(prs.slides)} slides."

# ──────────────────────────────────────────────────────────────
# Step 2: Agents Definition
# ──────────────────────────────────────────────────────────────

research_agent = create_react_agent(
    model=llm,
    tools=[research_tool],
    prompt=(
        "You are a research agent. Gather factual information about the topic. "
        "Summarize your findings logically so they can be easily converted into a presentation."
    )
)

# Note: The writer prompt is injected dynamically during execution to include User Branding
def get_writer_agent(branding: str):
    return create_react_agent(
        model=llm,
        tools=[],
        prompt=(
            "You are a presentation writer. Convert the provided research into a JSON array representing PPT slides.\n\n"
            f"BRANDING/TONE GUIDELINES: '{branding if branding else 'Professional and standard'}'\n\n"
            "RULES:\n"
            "- Output ONLY a valid JSON array.\n"
            "- First element is title slide: {\"title\": \"Title\", \"bullets\": [\"Subtitle\"], \"image_prompt\": \"\"}\n"
            "- Following 4-6 elements are content slides: {\"title\": \"Slide Title\", \"bullets\": [\"Point 1\", \"Point 2\"], \"image_prompt\": \"AI image generation prompt describing a relevant visual\"}\n"
            "- Bullet points MUST be concise (max 10 words each).\n"
            "- `image_prompt` MUST be a visual description (e.g., 'Modern corporate office, 8k, photorealistic'). "
            "Do NOT ask for text in the image. Keep it purely visual."
        )
    )

ppt_agent = create_react_agent(
    model=llm,
    tools=[create_ppt],
    prompt=(
        "You are a PowerPoint creation agent. Call the create_ppt tool using the JSON string provided by the writer. "
        "If it fails, ask for a correctly formatted JSON string."
    )
)

# ──────────────────────────────────────────────────────────────
# Step 3: Orchestration
# ──────────────────────────────────────────────────────────────

def run_agents(topic: str):
    print(f"\n{'='*60}")
    print(f"  Generating Presentation for: {topic}")
    print(f"{'='*60}")

    # Stage 1: Research
    print("\n[1/3] 🔍 Researching Topic...")
    research_result = research_agent.invoke({"messages": [{"role": "user", "content": topic}]})
    research_output = research_result["messages"][-1].content

    # Stage 2: Writer 
    print("\n[2/3] ✍️ Structuring Slides & Formatting Prompts...")
    writer_agent = get_writer_agent(USER_CONFIG["branding"])
    writer_prompt = f"Convert the following research into PPT JSON structure:\n\n{research_output}"
    
    writer_result = writer_agent.invoke({"messages": [{"role": "user", "content": writer_prompt}]})
    slides_json = writer_result["messages"][-1].content
    
    # Stage 3: PPT Agent creates file and fetches images
    print("\n[3/3] 🎨 Generating Images & Creating PowerPoint File...")
    ppt_result = ppt_agent.invoke({"messages": [{"role": "user", "content": slides_json}]})
    ppt_output = ppt_result["messages"][-1].content
    
    print("\n✅ Final Agent Output:")
    print(ppt_output)

# ──────────────────────────────────────────────────────────────
# Step 4: User Interaction Entry Point
# ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print(r"""
  __  __       _ _   _              _                  _       ____  ____ _____ 
 |  \/  |     | | | (_)            / \   __ _  ___ _ __| |_    |  _ \|  _ \_   _|
 | |\/| |_   _| | |_| |  _____    / _ \ / _` |/ _ \ '_ \ __|   | |_) | |_) || |  
 | |  | | |_| | |  _| | |_____|  / ___ \ (_| |  __/ | | | |_   |  __/|  __/ | |  
 |_|  |_|\__,_|_|\__|_|         /_/   \_\__, |\___|_| |_|\__|  |_|   |_|    |_|  
                                        |___/                                    
    """)
    
    # 1. Ask for Topic
    user_topic = input("\n👉 Enter the topic for your presentation: ").strip()
    if not user_topic:
        user_topic = "The Future of Artificial Intelligence"
        print(f"No topic provided. Defaulting to: {user_topic}")
    
    # 2. Ask for Template
    print("\n👉 Do you have a custom PowerPoint template? ")
    user_template = input("   (Enter the path to the .pptx file, or leave blank to use AI-generated design): ").strip()
    if user_template and os.path.exists(user_template):
        USER_CONFIG["template_path"] = user_template
        print(f"   [+] Valid template found: {user_template}")
    elif user_template:
        print("   [!] Template file not found. Falling back to default AI design.")
    else:
        print("   [+] Using default AI layout design.")

    # 3. Ask for Branding
    print("\n👉 Enter any specific branding/style guidelines.")
    user_branding = input("   (e.g., 'Minimalist design, corporate tone', or leave blank for standard professional style): ").strip()